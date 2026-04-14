from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Iterable, List, Tuple

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from config import AppConfig, get_config
from models import DocumentUnit
from text_utils import normalise_whitespace


def _bytes_limit_ok(data: bytes, cfg: AppConfig) -> bool:
    return len(data) <= cfg.MAX_FILE_MB * 1024 * 1024


def ingest_file_bytes(filename: str, data: bytes, cfg: AppConfig | None = None) -> List[DocumentUnit]:
    """
    Ingest a single file into provenance-preserving units.

    Report mapping:
    - Preserves PDF page numbers, DOCX paragraph indices, PPTX slide indices,
      CSV row indices, and TXT/MD line ranges.
    """
    cfg = cfg or get_config()
    if not _bytes_limit_ok(data, cfg):
        return []

    ext = Path(filename).suffix.lower().lstrip(".")
    units: List[DocumentUnit] = []
    unit_id_base = 0  # caller should reassign unit_ids if combining files

    if ext == "pdf":
        reader = PdfReader(io.BytesIO(data))
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            text = normalise_whitespace(text)
            if not text:
                continue
            units.append(
                DocumentUnit(
                    unit_id=unit_id_base + len(units),
                    text=text,
                    source=filename,
                    file_type="pdf",
                    location_type="page",
                    location=f"page {i}",
                    location_index=i,
                )
            )
        return units

    if ext == "docx":
        doc = DocxDocument(io.BytesIO(data))
        para_idx = 0
        for p in doc.paragraphs:
            t = normalise_whitespace(p.text or "")
            if not t:
                continue
            para_idx += 1
            units.append(
                DocumentUnit(
                    unit_id=unit_id_base + len(units),
                    text=t,
                    source=filename,
                    file_type="docx",
                    location_type="paragraph",
                    location=f"paragraph {para_idx}",
                    location_index=para_idx,
                )
            )
        return units

    if ext in {"txt", "md"}:
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("latin-1", errors="ignore")

        lines = [ln.rstrip("\n") for ln in text.splitlines()]
        batch_size = 40  
        start = 0
        while start < len(lines):
            end = min(start + batch_size, len(lines))
            block = "\n".join(lines[start:end]).strip()
            block = normalise_whitespace(block)
            if block:
                units.append(
                    DocumentUnit(
                        unit_id=unit_id_base + len(units),
                        text=block,
                        source=filename,
                        file_type=ext,
                        location_type="lines",
                        location=f"lines {start+1}-{end}",
                        location_index=start + 1,
                    )
                )
            start = end
        return units

    if ext == "csv":
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("latin-1", errors="ignore")

        reader = csv.reader(io.StringIO(text))
        for i, row in enumerate(reader, start=1):
            joined = " | ".join([c.strip() for c in row if c and c.strip()])
            joined = normalise_whitespace(joined)
            if not joined:
                continue
            units.append(
                DocumentUnit(
                    unit_id=unit_id_base + len(units),
                    text=joined,
                    source=filename,
                    file_type="csv",
                    location_type="row",
                    location=f"row {i}",
                    location_index=i,
                )
            )
        return units

    if ext == "pptx":
        prs = Presentation(io.BytesIO(data))
        for s_idx, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    t = shape.text_frame.text or ""
                    t = normalise_whitespace(t)
                    if t:
                        texts.append(t)
            slide_text = normalise_whitespace("\n".join(texts))
            if not slide_text:
                continue
            units.append(
                DocumentUnit(
                    unit_id=unit_id_base + len(units),
                    text=slide_text,
                    source=filename,
                    file_type="pptx",
                    location_type="slide",
                    location=f"slide {s_idx}",
                    location_index=s_idx,
                )
            )
        return units

    # Unknown extension: try to treat as text (best-effort).
    try:
        text = data.decode("utf-8")
    except Exception:
        return []
    text = normalise_whitespace(text)
    if not text:
        return []

    return [
        DocumentUnit(
            unit_id=unit_id_base,
            text=text,
            source=filename,
            file_type="unknown",
            location_type="file",
            location="whole file",
            location_index=None,
        )
    ]


def ingest_files(items: Iterable[Tuple[str, bytes]], cfg: AppConfig | None = None) -> List[DocumentUnit]:
    """
    Ingest multiple (filename, bytes) entries and renumber unit_ids globally.
    """
    cfg = cfg or get_config()
    all_units: List[DocumentUnit] = []
    next_id = 0

    for fname, data in items:
        units = ingest_file_bytes(fname, data, cfg=cfg)
        for u in units:
            all_units.append(
                DocumentUnit(
                    unit_id=next_id,
                    text=u.text,
                    source=u.source,
                    file_type=u.file_type,
                    location_type=u.location_type,
                    location=u.location,
                    location_index=u.location_index,
                )
            )
            next_id += 1

    return all_units


def ingest_uploaded_files(uploaded_files, cfg: AppConfig | None = None) -> List[DocumentUnit]:
    """
    Streamlit helper: accept a list of UploadedFile objects.
    """
    cfg = cfg or get_config()
    items: List[Tuple[str, bytes]] = []
    for f in uploaded_files:
        items.append((f.name, f.getvalue()))
    return ingest_files(items, cfg=cfg)